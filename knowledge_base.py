"""
知识库
"""
import os
import io
import hashlib
import pymysql
import config_data as config
from langchain_chroma import Chroma
from langchain_community.embeddings import DashScopeEmbeddings
from langchain_experimental.text_splitter import SemanticChunker
from langchain_text_splitters import RecursiveCharacterTextSplitter, MarkdownHeaderTextSplitter, Language
from datetime import datetime

# 修复 Python 3.12 + multiprocess 的 ResourceTracker shutdown 报错
import multiprocess.resource_tracker


def _noop_stop_locked(self):
    pass


multiprocess.resource_tracker.ResourceTracker._stop_locked = _noop_stop_locked

from sentence_transformers import SentenceTransformer
from PIL import Image


def _get_mysql_connection():
    return pymysql.connect(
        host=config.mysql_host,
        port=config.mysql_port,
        user=config.mysql_user,
        password=config.mysql_password,
        database=config.mysql_database,
        charset=config.mysql_charset,
    )


def _ensure_table():
    conn = _get_mysql_connection()
    try:
        with conn.cursor() as cursor:
            cursor.execute(
                """
                CREATE TABLE IF NOT EXISTS knowledge_md5 (
                    id INT AUTO_INCREMENT PRIMARY KEY,
                    md5_value VARCHAR(64) NOT NULL UNIQUE,
                    filename VARCHAR(255) DEFAULT '',
                    create_time DATETIME DEFAULT CURRENT_TIMESTAMP
                ) ENGINE=InnoDB DEFAULT CHARSET=utf8mb4
                """
            )
        conn.commit()
    finally:
        conn.close()


def get_string_md5(input_str: str, encoding='utf-8'):
    """将传入的字符串转换为md5字符串"""
    str_bytes = input_str.encode(encoding=encoding)
    md5_obj = hashlib.md5()
    md5_obj.update(str_bytes)
    return md5_obj.hexdigest()


def get_bytes_md5(data: bytes):
    """计算字节数据的MD5值"""
    md5_obj = hashlib.md5()
    md5_obj.update(data)
    return md5_obj.hexdigest()


def save_md5(md5_str: str, filename: str = "") -> int | None:
    """将传入的md5字符串，记录到MySQL数据库保存，返回 knowledge_md5 的 id"""
    _ensure_table()
    conn = _get_mysql_connection()
    try:
        with conn.cursor() as cursor:
            cursor.execute(
                "INSERT IGNORE INTO knowledge_md5 (md5_value, filename) VALUES (%s, %s)",
                (md5_str, filename),
            )
            conn.commit()
            cursor.execute("SELECT id FROM knowledge_md5 WHERE md5_value = %s", (md5_str,))
            row = cursor.fetchone()
            return row[0] if row else None
    finally:
        conn.close()


def check_md5(md5_str: str):
    """检查传入的md5字符串是否已经被处理过了
        return False(md5未处理过)  True(已经处理过，已有记录)
    """
    _ensure_table()
    conn = _get_mysql_connection()
    try:
        with conn.cursor() as cursor:
            cursor.execute(
                "SELECT 1 FROM knowledge_md5 WHERE md5_value = %s LIMIT 1",
                (md5_str,),
            )
            return cursor.fetchone() is not None
    finally:
        conn.close()


def delete_md5_by_filename(filename: str):
    """根据文件名删除 MySQL 中的 MD5 记录"""
    _ensure_table()
    conn = _get_mysql_connection()
    try:
        with conn.cursor() as cursor:
            cursor.execute(
                "DELETE FROM knowledge_md5 WHERE filename = %s",
                (filename,),
            )
        conn.commit()
    finally:
        conn.close()


def get_all_filenames() -> list[str]:
    """获取所有已入库的文件名"""
    _ensure_table()
    conn = _get_mysql_connection()
    try:
        with conn.cursor() as cursor:
            cursor.execute("SELECT DISTINCT filename FROM knowledge_md5 ORDER BY filename")
            return [row[0] for row in cursor.fetchall()]
    finally:
        conn.close()


def _cjk_ratio(text: str) -> float:
    if not text:
        return 0.0
    return sum(1 for ch in text if '一' <= ch <= '鿿') / max(len(text), 1)


def _printable_ratio(text: str) -> float:
    if not text:
        return 0.0
    printable = sum(1 for ch in text if ch.isprintable() or ch in '\n\r\t ')
    return printable / len(text)


def _is_scanned_pdf(file_bytes: bytes) -> bool:
    """快速探测 PDF 是否为扫描版（无文字层或文字层质量差）"""
    import fitz
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    probe_parts = []
    for page in doc:
        text = page.get_text()
        if text:
            probe_parts.append(text)
    doc.close()
    probe = '\n'.join(probe_parts)
    return (
        not probe.strip()
        or _cjk_ratio(probe) < 0.05
        or _printable_ratio(probe) < 0.8
    )


def _read_pdf_pages(file_bytes: bytes) -> list[tuple[int, str]]:
    """逐页读取 PDF 文字，返回 [(page_num, page_text), ...]
    文字版/扫描版均使用 TableTransformer 做表格识别
    """
    import fitz

    if _is_scanned_pdf(file_bytes):
        import sys
        from concurrent.futures import ThreadPoolExecutor, as_completed
        from ocr_pdf import _get_ocr, _bytes_to_array
        ocr = _get_ocr()
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        total = len(doc)

        # 阶段1: 渲染所有页面为图片（IO 密集，快速）
        page_images = []
        for page_num in range(total):
            pix = doc[page_num].get_pixmap(dpi=150)
            page_images.append((page_num + 1, _bytes_to_array(pix.tobytes("png"))))
        doc.close()

        # 阶段2: 多线程并行 OCR（CPU 密集，PaddleOCR 推理时释放 GIL）
        results = {}

        def _ocr_page(pn, img_array):
            result = list(ocr.ocr(img_array))
            text = ""
            if result and result[0]:
                rec_texts = result[0].get("rec_texts", [])
                text = "\n".join(rec_texts)
            return (pn, text)

        with ThreadPoolExecutor(max_workers=4) as pool:
            futures = {pool.submit(_ocr_page, pn, img): pn for pn, img in page_images}
            for future in as_completed(futures):
                pn, text = future.result()
                results[pn] = text
                sys.stderr.write(f"\r[OCR] {len(results)}/{total} 页   ")
                sys.stderr.flush()

        pages = [(pn, results[pn]) for pn in sorted(results)]
        print(file=sys.stderr)
    else:
        # 文字版：逐页提取文字层
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        pages = []
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text()
            pages.append((page_num + 1, text))
        doc.close()

    # 表格检测：使用 TableTransformer + PaddleOCR（文字版和扫描版通用）
    try:
        from table_extractor import extract_tables_from_pdf
        tables = extract_tables_from_pdf(file_bytes)
    except Exception:
        tables = []

    page_tables = {}
    for t in tables:
        pn = t["page_num"]
        page_tables.setdefault(pn, []).append(t)

    # 将表格 Markdown 追加到对应页面文本
    for i, (page_num, text) in enumerate(pages):
        for t in page_tables.get(page_num, []):
            text += f"\n[表格]\n{t['markdown']}\n"
        pages[i] = (page_num, text)

    return pages


def _strip_html_tags(text: str) -> str:
    """去除 HTML 标签（<img>、<div> 等），保留纯文本"""
    import re
    # 去除所有 HTML 标签
    text = re.sub(r'<[^>]+>', '', text)
    # 去除标签残留的多余空行
    text = re.sub(r'\n{3,}', '\n\n', text)
    return text.strip()


def read_file_content(file_bytes: bytes, filename: str) -> str:
    """根据文件扩展名读取文件内容为字符串"""
    ext = os.path.splitext(filename)[1].lower()

    if ext == '.txt':
        return file_bytes.decode('utf-8')

    elif ext == '.md':
        return _strip_html_tags(file_bytes.decode('utf-8'))

    elif ext == '.docx':
        from table_extractor import read_docx_with_tables
        return read_docx_with_tables(file_bytes)

    elif ext == '.xlsx':
        import io
        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(file_bytes), read_only=True)
        text_parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            text_parts.append(f'[Sheet: {sheet_name}]')
            rows = list(ws.iter_rows(values_only=True))
            non_empty = []
            for row in rows:
                row_str = [str(c) if c is not None else '' for c in row]
                if any(c.strip() for c in row_str):
                    non_empty.append(row_str)
            if non_empty:
                max_cols = max(len(r) for r in non_empty)
                for i, row in enumerate(non_empty):
                    while len(row) < max_cols:
                        row.append('')
                    text_parts.append('| ' + ' | '.join(row) + ' |')
                    if i == 0:
                        text_parts.append('| ' + ' | '.join(['---'] * max_cols) + ' |')
        wb.close()
        return '\n'.join(text_parts)

    elif ext == '.pptx':
        import io
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        text_parts = []
        for slide_idx, slide in enumerate(prs.slides, 1):
            text_parts.append(f'[幻灯片 {slide_idx}]')
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            text_parts.append(para.text)
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = '\t'.join([cell.text for cell in row.cells])
                        if row_text.strip():
                            text_parts.append(row_text)
        return '\n'.join(text_parts)

    elif ext == '.csv':
        import io
        import csv
        text_parts = []
        csv_reader = csv.reader(io.StringIO(file_bytes.decode('utf-8')))
        for row in csv_reader:
            row_text = '\t'.join([cell.strip() if cell else '' for cell in row])
            if row_text.strip():
                text_parts.append(row_text)
        return '\n'.join(text_parts)

    else:
        raise ValueError(f"不支持的文件格式: {ext}")


def extract_images_from_pdf(file_bytes: bytes, filename: str) -> list[dict]:
    """从 PDF 中提取图片，返回 [{page_num, image_bytes, ext}] 列表"""
    import fitz  # PyMuPDF
    images = []
    doc = fitz.open(stream=file_bytes, filetype="pdf")

    for page_num in range(len(doc)):
        page = doc[page_num]
        # 获取页面中的所有图片
        image_list = page.get_images(full=True)
        for img_idx, img_info in enumerate(image_list):
            xref = img_info[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            ext = base_image["ext"]  # png, jpeg, etc.

            # 过滤太小的图片（可能是图标、装饰元素）
            if len(image_bytes) < 2048:  # < 2KB
                continue

            images.append({
                "page_num": page_num + 1,
                "image_bytes": image_bytes,
                "ext": ext,
                "filename": filename,
            })

    doc.close()
    return images


_clip_model = None


def _get_clip_model():
    """延迟加载 CLIP 模型单例（首次调用时自动下载约 600MB）"""
    global _clip_model
    if _clip_model is None:
        _clip_model = SentenceTransformer(config.image_model_path, local_files_only=True)
        print("OK")
    return _clip_model


def _embed_images_batch(image_items: list[dict], batch_size: int = 16) -> list[list[float]]:
    """批量将图片向量化（本地 CLIP 模型），返回 embedding 列表"""
    if not image_items:
        return []

    model = _get_clip_model()
    all_embeddings = []
    for start in range(0, len(image_items), batch_size):
        batch = image_items[start:start + batch_size]
        pil_images = [Image.open(io.BytesIO(item["image_bytes"])) for item in batch]
        embeddings = model.encode(pil_images, batch_size=batch_size, show_progress_bar=False)
        all_embeddings.extend(embeddings.tolist())

    return all_embeddings


def _embed_text_for_image_search(query: str) -> list[float]:
    """将查询文本向量化到图片向量空间（本地 CLIP 模型），用于图文检索"""
    model = _get_clip_model()
    embedding = model.encode([query], show_progress_bar=False)
    return embedding[0].tolist()


def clean_text(text: str) -> str:
    """数据清洗：去页眉页脚页码、清洗空白、删除目录/版权页、过滤短句碎片"""
    import re

    lines = text.split('\n')

    # ========== 1. 页眉页脚 + 页码删除 ==========
    # 1a. 删除独立的页码行（纯数字）
    lines = [l for l in lines if not re.match(r'^\s*\d{1,4}\s*$', l)]

    # 1b. 删除常见的页眉页脚模式
    header_footer_patterns = [
        r'^第[一二三四五六七八九十百千\d]+[章节].{0,10}$',  # "第X章" 单行
        r'^\d{4}[\/\-\.]\d{1,2}[\/\-\.]\d{1,2}$',         # 日期
        r'^[\(（]?\d{4}[\)）]?$',                            # 年份括号
    ]
    cleaned_lines = []
    for line in lines:
        is_header_footer = any(re.match(p, line.strip()) for p in header_footer_patterns)
        if not is_header_footer:
            cleaned_lines.append(line)
    lines = cleaned_lines

    # ========== 2. 空白/换行符清洗 ==========
    cleaned_lines = []
    for line in lines:
        stripped = line.strip()
        if stripped:
            cleaned_lines.append(stripped)
        else:
            # 保留单个空行作为段落分隔，但折叠连续空行
            if cleaned_lines and cleaned_lines[-1] != '':
                cleaned_lines.append('')
    lines = cleaned_lines

    if not lines:
        return ''

    # 合并同一段落内被 PDF 断开的行（行末没有句末标点则与下一行合并）
    merged_lines = []
    sentence_endings = set('。！？…"」』》）.!?"')
    buffer = ''
    for line in lines:
        if not line:
            if buffer:
                merged_lines.append(buffer)
                buffer = ''
            merged_lines.append('')
            continue
        if buffer:
            # 如果上一行末尾不是句末标点，说明是被 PDF 断开的同一句话
            if buffer[-1] not in sentence_endings and len(buffer) < config.max_chunk_api_limit:
                buffer += line
            else:
                merged_lines.append(buffer)
                buffer = line
        else:
            buffer = line
    if buffer:
        merged_lines.append(buffer)
    lines = merged_lines

    # 再次归一化空白行
    cleaned_lines = []
    for line in lines:
        if line:
            # 行内多余空白归一
            cleaned_lines.append(re.sub(r'\s+', ' ', line).strip())
        else:
            if cleaned_lines and cleaned_lines[-1] != '':
                cleaned_lines.append('')
    lines = [l for l in cleaned_lines if l]
    if not lines:
        return ''

    # ========== 3. 目录/版权页检测 ==========
    # 3a. 目录页检测：连续多行命中 toc 模式的行视为目录区
    toc_scores = []
    for line in lines:
        score = 0
        for pattern in config.toc_line_patterns:
            if re.search(pattern, line):
                score += 1
                if score >= 1:
                    break
        toc_scores.append(score)

    # 标记连续目录区（5行以内有3行命中，视为目录块起点和终点）
    in_toc_region = [False] * len(lines)
    i = 0
    while i < len(lines):
        window_end = min(i + 10, len(lines))
        window_scores = sum(toc_scores[i:window_end])
        if window_scores >= 3:
            # 找到连续命中区域
            j = i
            while j < len(lines) and toc_scores[j] > 0:
                j += 1
            # 如果连续命中行 >= 5，标记整个区域
            consecutive_hits = sum(1 for k in range(i, min(j, len(lines))) if toc_scores[k] > 0)
            if consecutive_hits >= 5:
                for k in range(i, min(j, len(lines))):
                    in_toc_region[k] = True
            i = j
        else:
            i += 1

    lines = [l for idx, l in enumerate(lines) if not in_toc_region[idx]]

    # 3b. 版权页检测：包含版权关键词的段落（前后各2行一起删除）
    copyright_set = set()
    for idx, line in enumerate(lines):
        for kw in config.copyright_keywords:
            if kw in line:
                for offset in range(-2, 3):
                    target = idx + offset
                    if 0 <= target < len(lines):
                        copyright_set.add(target)
                break

    lines = [l for idx, l in enumerate(lines) if idx not in copyright_set]

    # ========== 4. 短句碎片过滤 ==========
    lines = [l for l in lines if len(l) >= config.min_line_char_length]

    return '\n'.join(lines)


def _split_markdown_by_headers(text: str) -> list[str]:
    """基于 Markdown 标题层级（# ## ###）做结构化切分，保留段落完整性"""
    splitter = MarkdownHeaderTextSplitter(
        headers_to_split_on=[("#", "h1"), ("##", "h2"), ("###", "h3")],
        strip_headers=False,
    )
    docs = splitter.split_text(text)
    chunks = []
    for doc in docs:
        header_parts = []
        for key in ("h1", "h2", "h3"):
            val = doc.metadata.get(key)
            if val:
                header_parts.append(val)
        if header_parts:
            chunks.append(" > ".join(header_parts) + "\n" + doc.page_content)
        else:
            chunks.append(doc.page_content)
    return chunks


def _split_by_slides(text: str) -> list[str]:
    """按幻灯片边界 [幻灯片 N] 切分"""
    import re
    chunks = re.split(r'\n(?=\[幻灯片 \d+\])', text)
    return [c.strip() for c in chunks if c.strip()]


def _split_by_sheets(text: str) -> list[str]:
    """按 Sheet 边界 [Sheet: ...] 切分"""
    import re
    chunks = re.split(r'\n(?=\[Sheet: )', text)
    return [c.strip() for c in chunks if c.strip()]


class KnowledgeBaseService(object):
    def __init__(self):
        os.makedirs(config.persist_directory, exist_ok=True)

        self.chroma = Chroma(
            collection_name=config.collection_name,
            embedding_function=DashScopeEmbeddings(
                dashscope_api_key=config.dashscope_api_key,
                model=config.embedding_model_name,
            ),
            persist_directory=config.persist_directory,
        )

        # 图片专用向量库
        self.image_chroma = Chroma(
            collection_name=config.collection_name + "_images",
            embedding_function=None,  # 手动传入 embedding
            persist_directory=config.persist_directory,
        )

        self.spliter = SemanticChunker(
            embeddings=DashScopeEmbeddings(
                dashscope_api_key=config.dashscope_api_key,
                model=config.embedding_model_name,
            ),
            breakpoint_threshold_type=config.semantic_breakpoint_threshold_type,
            breakpoint_threshold_amount=config.semantic_breakpoint_threshold_amount,
        )

    def upload_by_str(self, data: str, filename: str, page_num: int = None,
                      page_ranges: list[tuple] = None, file_ext: str = None):
        """将传入的字符串，进行向量化，存入向量数据库中

        Args:
            page_ranges: [(page_num, start, end), ...] 每页在已清洗文本中的字符偏移，
                         用于跨页面语义切分后将 chunk 映射回原始页码（用于图片检索）。
        """
        if page_ranges is None:
            # 结构化格式（有自身标记语法）跳过清洗，避免破坏标题/表格/Slide 标记
            if file_ext in ('.md', '.docx', '.pptx', '.xlsx'):
                if not data or not data.strip():
                    return "[跳过]内容为空", None
            else:
                data = clean_text(data)
                if not data or not data.strip():
                    return "[跳过]清洗后内容为空", None
        else:
            # PDF 路径：数据已经逐页清洗过，直接使用
            if not data or not data.strip():
                return "[跳过]清洗后内容为空", None

        md5_hex = get_string_md5(data)

        if check_md5(md5_hex):
            return "[跳过]内容已经存在知识库中", None

        # 兜底切分器：处理超长块的二次切分
        _fallback = RecursiveCharacterTextSplitter(
            chunk_size=config.max_chunk_api_limit,
            chunk_overlap=config.max_chunk_overlap,
            separators=["\n\n", "\n", "。", "！", "？", ".", "!", "?", " ", ""],
        )

        if len(data) > config.max_chunk_api_limit:
            if file_ext in ('.md', '.docx'):
                # 结构化文档：基于标题层级切分（md 用 # 语法，docx 已转为 # 格式）
                knowledge_chunks = _split_markdown_by_headers(data)
            elif file_ext == '.pptx':
                # PPT：按幻灯片边界切分
                knowledge_chunks = _split_by_slides(data)
            elif file_ext == '.xlsx':
                # Excel：按 Sheet 切分
                knowledge_chunks = _split_by_sheets(data)
            elif file_ext == '.pdf':
                # PDF：先预切分到安全范围，再逐段语义切分细化
                pre_chunks = _fallback.split_text(data)
                knowledge_chunks = []
                semantic_errors = 0
                for pre in pre_chunks:
                    if len(pre) <= config.max_chunk_api_limit:
                        try:
                            sub = self.spliter.split_text(pre)
                            knowledge_chunks.extend(sub)
                        except Exception:
                            semantic_errors += 1
                            knowledge_chunks.append(pre)
                    else:
                        knowledge_chunks.append(pre)
                if semantic_errors > 0:
                    import sys
                    sys.stderr.write(
                        f"[SemanticChunker] {semantic_errors}/{len(pre_chunks)} 段语义切分失败，"
                        f"已回退为递归切分结果\n"
                    )
                    sys.stderr.flush()
            else:
                # txt 等：语义切分，失败回退递归切分
                try:
                    knowledge_chunks = self.spliter.split_text(data)
                except Exception:
                    knowledge_chunks = _fallback.split_text(data)

            # 兜底过滤空块 + 超长块二次切分
            min_len = 100 if file_ext == '.md' else 0
            safe_chunks = []
            for chunk in knowledge_chunks:
                chunk = chunk.strip()
                if not chunk or len(chunk) < min_len:
                    continue
                if len(chunk) > config.max_chunk_api_limit:
                    safe_chunks.extend(_fallback.split_text(chunk))
                else:
                    safe_chunks.append(chunk)
            knowledge_chunks = safe_chunks
        else:
            knowledge_chunks = [data]

        # 将 chunk 映射回页码
        chunk_page_nums = self._map_chunks_to_pages(data, knowledge_chunks, page_ranges)

        metadatas = []
        for i, chunk in enumerate(knowledge_chunks):
            meta = {
                "source": filename,
                "create_time": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                "operator": "小黄",
            }
            pns = chunk_page_nums[i] if i < len(chunk_page_nums) else []
            if pns:
                meta["page_nums"] = ",".join(str(p) for p in pns)
            elif page_num is not None:
                meta["page_nums"] = str(page_num)
            metadatas.append(meta)

        try:
            self.chroma.add_texts(knowledge_chunks, metadatas=metadatas)
        except Exception:
            # Chroma 入库失败，逐条重试跳过坏块
            for chunk, meta in zip(knowledge_chunks, metadatas):
                try:
                    self.chroma.add_texts([chunk], metadatas=[meta])
                except Exception:
                    pass

        save_md5(md5_hex, filename)

        count = len(knowledge_chunks)
        return f"[成功]内容已经成功载入向量库（{count} 个文本块）", md5_hex

    @staticmethod
    def _map_chunks_to_pages(text: str, chunks: list[str],
                             page_ranges: list[tuple]) -> list[list[int]]:
        """将每个 chunk 映射到它所属的页码列表（跨页 chunk 会包含多页）"""
        if not page_ranges:
            return [[] for _ in chunks]

        result = []
        for chunk in chunks:
            idx = text.find(chunk)
            if idx == -1:
                result.append([])
                continue
            chunk_start = idx
            chunk_end = idx + len(chunk)
            pages = []
            for pn, p_start, p_end in page_ranges:
                if chunk_start < p_end and chunk_end > p_start:
                    pages.append(pn)
            result.append(pages)
        return result

    def upload_by_bytes(self, file_bytes: bytes, filename: str):
        """将传入的文件字节数据解析后向量化存入数据库（包含文字和图片）"""
        ext = os.path.splitext(filename)[1].lower()
        file_type = ext.lstrip('.')
        md5_value = None

        try:
            # 1. 处理文本：PDF 先逐页清洗再合并做跨页面语义切分，其他格式全量处理
            if ext == '.pdf':
                pages = _read_pdf_pages(file_bytes)
                cleaned_pages = []
                for page_num, page_text in pages:
                    cleaned = clean_text(page_text)
                    if cleaned and cleaned.strip():
                        cleaned_pages.append((page_num, cleaned))

                if not cleaned_pages:
                    result = "[跳过]所有页面清洗后内容为空"
                else:
                    merged_parts = []
                    page_ranges = []
                    offset = 0
                    for page_num, text in cleaned_pages:
                        merged_parts.append(text)
                        end = offset + len(text)
                        page_ranges.append((page_num, offset, end))
                        offset = end + 1

                    merged_text = "\n".join(merged_parts)
                    result, md5_value = self.upload_by_str(merged_text, filename, page_ranges=page_ranges, file_ext='.pdf')
            else:
                text = read_file_content(file_bytes, filename)
                result, md5_value = self.upload_by_str(text, filename, file_ext=ext)

            # 2. 处理图片（仅文字版 PDF 提取图片，扫描版跳过）
            if ext == '.pdf' and not _is_scanned_pdf(file_bytes):
                try:
                    images = extract_images_from_pdf(file_bytes, filename)
                    if images:
                        self.upload_images(images)
                        result += f"；已提取 {len(images)} 张图片"
                except ImportError as e:
                    result += f"；[跳过图片] {e}"
                except Exception as e:
                    result += f"；[图片处理异常] {e}"

            # 3. 记录上传日志
            from knowledge_manage.db import insert_upload_record
            if md5_value is not None:
                insert_upload_record(md5_value, filename, file_type, 'success')
            else:
                insert_upload_record(None, filename, file_type, 'failed')
            return result

        except Exception:
            from knowledge_manage.db import insert_upload_record
            insert_upload_record(None, filename, file_type, 'failed')
            raise

    def upload_images(self, images: list[dict]):
        """将图片列表向量化后存入图片向量库"""
        if not images:
            return

        # 批量向量化
        embeddings = _embed_images_batch(images)

        # 构造文档内容、metadata、ids
        texts = []
        metadatas = []
        ids = []
        for i, img in enumerate(images):
            texts.append(f"[图片] 来源: {img['filename']}, 第 {img['page_num']} 页")
            metadatas.append({
                "source": img["filename"],
                "page_num": img["page_num"],
                "type": "image",
                "ext": img["ext"],
                "create_time": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                "operator": "小黄",
            })
            ids.append(f"img_{img['filename']}_p{img['page_num']}_{i}")

        # 直接调用底层 Chroma collection 的 add 方法
        try:
            self.image_chroma._collection.add(
                ids=ids,
                embeddings=embeddings,
                documents=texts,
                metadatas=metadatas,
            )
        except ValueError as e:
            err_msg = str(e)
            if "dimension" in err_msg.lower():
                raise ValueError(
                    f"图片向量维度不匹配，请重置图片向量集合（当前模型与旧集合维度不一致）: {err_msg}"
                )
            raise

    def reset_image_collection(self):
        """删除并重建图片向量集合（用于切换 embedding 模型后维度不匹配的情况）"""
        try:
            self.image_chroma._client.delete_collection(
                config.collection_name + "_images"
            )
        except Exception:
            pass
        self.image_chroma = Chroma(
            collection_name=config.collection_name + "_images",
            embedding_function=None,
            persist_directory=config.persist_directory,
        )
        return "图片向量集合已重置，请重新上传含图片的 PDF"

    def delete_by_filename(self, filename: str):
        """根据文件名从向量库和 MySQL 中删除知识文档"""
        self.chroma.delete(where={"source": filename})
        # 同时删除图片向量库中该文件的数据
        try:
            self.image_chroma._collection.delete(where={"source": filename})
        except Exception:
            pass
        delete_md5_by_filename(filename)
        return f"[成功]已删除知识库文档: {filename}"

    def clear_all(self):
        """一键清空所有知识库文档（文本 + 图片 + MySQL 记录）"""
        # 清空文本向量集合
        all_ids = self.chroma.get()["ids"]
        if all_ids:
            self.chroma.delete(ids=all_ids)
        # 清空图片向量集合
        try:
            img_ids = self.image_chroma._collection.get()["ids"]
            if img_ids:
                self.image_chroma._collection.delete(ids=img_ids)
        except Exception:
            pass
        # 清空 MySQL MD5 记录
        conn = _get_mysql_connection()
        try:
            with conn.cursor() as cursor:
                cursor.execute("DELETE FROM file_uploading")
                cursor.execute("DELETE FROM knowledge_md5")
            conn.commit()
        finally:
            conn.close()
        return "[成功]已清空所有知识库文档"

    def get_all_filenames(self) -> list[str]:
        """获取向量库中所有不重复的文件名"""
        results = self.chroma.get()
        sources = set()
        if results and results.get("metadatas"):
            for meta in results["metadatas"]:
                if meta and meta.get("source"):
                    sources.add(meta["source"])
        return sorted(list(sources))
    


